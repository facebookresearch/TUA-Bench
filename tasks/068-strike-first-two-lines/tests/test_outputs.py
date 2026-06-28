# Copyright (c) Meta Platforms, Inc. and affiliates.
# All rights reserved.
#
# This source code is licensed under the license found in the
# LICENSE file in the root directory of this source tree.

TASK_SPEC = {'artifact_paths': ['/app/New_Club_Spring_2018_Training.pptx'],
 'example_id': '550ce7e7-747b-495f-b122-acdc4d0b8e54',
 'inputs': [{'dest': 'New_Club_Spring_2018_Training.pptx',
             'url': 'https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/libreoffice_impress/550ce7e7-747b-495f-b122-acdc4d0b8e54/New_Club_Spring_2018_Training.data'}],
 'instruction': "I am checking our soccer club's to-do list for the last semester and adding "
                'strike-through sign on the line we have already accomplished. Could you help me '
                'add a strike-through on the first and second line?',
 'live_task': False,
 'metric_spec': {'conj': 'or',
                 'expected': [{'dest': 'New_Club_Spring_2018_Training_Gold.pptx',
                               'type': 'cloud_file',
                               'url': 'https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/libreoffice_impress/550ce7e7-747b-495f-b122-acdc4d0b8e54/New_Club_Spring_2018_Training_with_strike.data'},
                              {'dest': 'New_Club_Spring_2018_Training_Gold_all_fonts_1.pptx',
                               'type': 'cloud_file',
                               'url': 'https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/libreoffice_impress/550ce7e7-747b-495f-b122-acdc4d0b8e54/New_Club_Spring_2018_Training_Gold_all_fonts_1.pptx'}],
                 'func': ['compare_pptx_files', 'compare_pptx_files'],
                 'options': [{'enable_debug': False}, {'enable_debug': False}],
                 'result': ['/app/New_Club_Spring_2018_Training.pptx',
                            '/app/New_Club_Spring_2018_Training.pptx']},
 'oracle': {'mode': 'fetch_copy',
            'url': 'https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/libreoffice_impress/550ce7e7-747b-495f-b122-acdc4d0b8e54/New_Club_Spring_2018_Training_with_strike.data'},
 'output_path': '/app/New_Club_Spring_2018_Training.pptx',
 'slug': '068-strike-first-two-lines'}

import json
import logging
import os
import shutil
import subprocess
import tempfile
import xml.etree.ElementTree as ET
import zipfile
from math import sqrt
from pathlib import Path
from typing import Dict, List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

try:
    import librosa
except ImportError:
    librosa = None

try:
    import numpy as np
except ImportError:
    np = None

try:
    from PIL import Image
except ImportError:
    Image = None

try:
    from fastdtw import fastdtw
except ImportError:
    fastdtw = None

try:
    from scipy.spatial.distance import cosine
except ImportError:
    cosine = None

try:
    from skimage.metrics import structural_similarity as ssim
except ImportError:
    ssim = None

logger = logging.getLogger("tua.impress.osworld")
logging.basicConfig(level=logging.INFO)

CACHE_DIR = Path("/tmp") / TASK_SPEC["slug"]


def get_app_dir() -> Path:
    return Path(os.environ.get("APP_DIR", "/app"))


def get_agent_home() -> Path:
    return Path(os.environ.get("AGENT_HOME", "/home/agent"))


def localize_path(path: str) -> Path:
    if path.startswith("/app/"):
        return get_app_dir() / path.removeprefix("/app/")
    if path.startswith("/home/agent/"):
        return get_agent_home() / path.removeprefix("/home/agent/")
    return Path(path)


def _require_optional_dependencies(**modules) -> None:
    missing = [name for name, module in modules.items() if module is None]
    if missing:
        raise RuntimeError(f"Missing optional verifier dependencies: {', '.join(missing)}")


def fetch_cloud_file(url: str, dest: str) -> Path:
    CACHE_DIR.mkdir(parents=True, exist_ok=True)
    target = CACHE_DIR / dest
    if not target.exists():
        subprocess.run(
            ["curl", "-fsSL", "--retry", "10", "--retry-all-errors", "--retry-delay", "2", url, "-o", str(target)],
            check=True,
        )
    return target


def materialize_inputs(app_dir: Path) -> None:
    app_dir.mkdir(parents=True, exist_ok=True)
    get_agent_home().mkdir(parents=True, exist_ok=True)
    for item in TASK_SPEC.get("inputs", []):
        destination = app_dir / item["dest"]
        destination.parent.mkdir(parents=True, exist_ok=True)
        if not destination.exists():
            shutil.copyfile(fetch_cloud_file(item["url"], item["dest"]), destination)


def run_oracle(app_dir: Path, agent_home: Path) -> None:
    env = os.environ.copy()
    env["APP_DIR"] = str(app_dir)
    env["AGENT_HOME"] = str(agent_home)
    subprocess.run(
        ["bash", str(Path(__file__).resolve().parents[1] / "solution" / "solve.sh")],
        check=True,
        env=env,
    )


def check_presenter_console_disable(config_file_path):
    try:
        tree = ET.parse(config_file_path)
        root = tree.getroot()

        namespaces = {
            'oor': 'http://openoffice.org/2001/registry'
        }

        for item in root.findall(
                ".//item[@oor:path='/org.openoffice.Office.Impress/Misc/Start']/prop[@oor:name='EnablePresenterScreen']",
                namespaces):
            presenter_screen_enabled = item.find('value').text
            if presenter_screen_enabled.lower() == 'false':
                return 1.
            else:
                return 0.
        return 0.
    except Exception as e:
        logger.error(f"Error: {e}")
        return 0.


def check_image_stretch_and_center(modified_ppt, original_ppt):
    original_pres = Presentation(original_ppt)
    modified_pres = Presentation(modified_ppt)

    original_slide = original_pres.slides[0]
    modified_slide = modified_pres.slides[0]

    original_slide_images = [shape for shape in original_slide.shapes if shape.shape_type == 13]
    modified_slide_images = [shape for shape in modified_slide.shapes if shape.shape_type == 13]

    if not original_slide_images:
        return 0.

    the_image = original_slide_images[0]

    the_modified_image = None
    for modified_image in modified_slide_images:
        if the_image.image.blob == modified_image.image.blob:
            the_modified_image = modified_image

    if the_modified_image is None:
        return 0.

    if (abs(the_modified_image.width - original_pres.slide_width) > Inches(0.5) or
            abs(the_modified_image.height - original_pres.slide_height) > Inches(0.5) or
            abs(the_modified_image.left - (original_pres.slide_width - the_modified_image.width) / 2) > Inches(0.5) or
            abs(the_modified_image.top - (original_pres.slide_height - the_modified_image.height) / 2) > Inches(0.5)):
        return 0.

    return 1.


def get_all_text_shapes(slide):
    def extract_text_shapes(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, "shapes"):
            for sub_shape in shape.shapes:
                results.extend(extract_text_shapes(sub_shape))
        return results

    all_text_shapes = []
    for shape in slide.shapes:
        all_text_shapes.extend(extract_text_shapes(shape))
    return all_text_shapes


def compare_pptx_files(file1_path, file2_path, **options):
    prs1 = Presentation(file1_path)
    prs2 = Presentation(file2_path)

    approximately_tolerance = options.get("approximately_tolerance", 0.005)

    def is_approximately_equal(val1, val2, tolerance=approximately_tolerance):
        if val1 == val2:
            return True
        if val1 == 0 and val2 == 0:
            return True
        if val1 == 0 or val2 == 0:
            return False
        return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance

    def nonempty_runs(para):
        return [r for r in para.runs if (r.text or "").strip() != ""]

    examine_number_of_slides = options.get("examine_number_of_slides", True)
    examine_shape = options.get("examine_shape", True)
    examine_text = options.get("examine_text", True)
    examine_indent = options.get("examine_indent", True)
    examine_font_name = options.get("examine_font_name", True)
    examine_font_size = options.get("examine_font_size", True)
    examine_font_bold = options.get("examine_font_bold", True)
    examine_font_italic = options.get("examine_font_italic", True)
    examine_color_rgb = options.get("examine_color_rgb", True)
    examine_font_underline = options.get("examine_font_underline", True)
    examine_strike_through = options.get("examine_strike_through", True)
    examine_alignment = options.get("examine_alignment", True)
    examine_title_bottom_position = options.get("examine_title_bottom_position", False)
    examine_table_bottom_position = options.get("examine_table_bottom_position", False)
    examine_run_count = options.get("examine_run_count", True)
    examine_right_position = options.get("examine_right_position", False)
    examine_top_position = options.get("examine_top_position", False)
    examine_shape_for_shift_size = options.get("examine_shape_for_shift_size", False)
    examine_image_size = options.get("examine_image_size", False)
    examine_modify_height = options.get("examine_modify_height", False)
    examine_bullets = options.get("examine_bullets", True)
    examine_background_color = options.get("examine_background_color", True)
    examine_note = options.get("examine_note", True)

    if len(prs1.slides) != len(prs2.slides) and examine_number_of_slides:
        return 0

    slide_idx = 0
    for slide1, slide2 in zip(prs1.slides, prs2.slides):
        slide_idx += 1

        def get_slide_background_color(slide):
            fill = slide.background.fill
            if fill.type == 1:
                return fill.fore_color.rgb
            elif fill.type == 5:
                master_fill = slide.slide_layout.slide_master.background.fill
                if master_fill.type == 1:
                    return master_fill.fore_color.rgb
                else:
                    return None
            else:
                return None

        if get_slide_background_color(slide1) != get_slide_background_color(slide2) and examine_background_color:
            return 0

        def get_slide_notes(slide):
            notes_slide = slide.notes_slide
            if notes_slide:
                return notes_slide.notes_text_frame.text
            else:
                return None

        if get_slide_notes(slide1).strip() != get_slide_notes(slide2).strip() and examine_note:
            return 0

        text_shapes1 = get_all_text_shapes(slide1)
        text_shapes2 = get_all_text_shapes(slide2)

        if len(slide1.shapes) != len(slide2.shapes):
            return 0

        for shape1, shape2 in zip(slide1.shapes, slide2.shapes):
            if examine_title_bottom_position:
                if hasattr(shape1, "text") and hasattr(shape2, "text") and shape1.text == shape2.text:
                    if shape1.text == "Product Comparison" and (shape1.top <= shape2.top or shape1.top < 3600000):
                        return 0
                elif (not is_approximately_equal(shape1.left, shape2.left) or
                      not is_approximately_equal(shape1.top, shape2.top) or
                      not is_approximately_equal(shape1.width, shape2.width) or
                      not is_approximately_equal(shape1.height, shape2.height)):
                    return 0

            if examine_table_bottom_position:
                if slide_idx == 3 and shape1.shape_type == 19 and shape2.shape_type == 19:
                    if shape1.top <= shape2.top or shape1.top < 3600000:
                        return 0
                elif (not is_approximately_equal(shape1.left, shape2.left) or
                      not is_approximately_equal(shape1.top, shape2.top) or
                      not is_approximately_equal(shape1.width, shape2.width) or
                      not is_approximately_equal(shape1.height, shape2.height)):
                    return 0

            if examine_right_position:
                if slide_idx == 2 and not hasattr(shape1, "text") and not hasattr(shape2, "text"):
                    if shape1.left <= shape2.left or shape1.left < 4320000:
                        return 0

            if examine_top_position:
                if slide_idx == 2 and shape1.shape_type == 13 and shape2.shape_type == 13:
                    if shape1.top >= shape2.top or shape1.top > 1980000:
                        return 0

            if examine_shape_for_shift_size:
                if (not is_approximately_equal(shape1.left, shape2.left) or
                    not is_approximately_equal(shape1.top, shape2.top) or
                    not is_approximately_equal(shape1.width, shape2.width) or
                    not is_approximately_equal(shape1.height, shape2.height)):
                    if not (hasattr(shape1, "text") and hasattr(shape2, "text")
                            and shape1.text == shape2.text and shape1.text == "Elaborate on what you want to discuss."):
                        return 0

            both_empty_text = (
                (not getattr(shape1, "text", None) or (shape1.text or "").strip() == "") and
                (not getattr(shape2, "text", None) or (shape2.text or "").strip() == "")
            )

            if (
                    not is_approximately_equal(shape1.left, shape2.left) or
                    not is_approximately_equal(shape1.top, shape2.top) or
                    not is_approximately_equal(shape1.width, shape2.width) or
                    not is_approximately_equal(shape1.height, shape2.height)) and examine_shape and not both_empty_text:
                return 0

            if examine_image_size:
                if shape1.shape_type == 13 and shape2.shape_type == 13:
                    if not is_approximately_equal(shape1.width, shape2.width) or not is_approximately_equal(shape1.height, shape2.height):
                        return 0
                elif (not is_approximately_equal(shape1.left, shape2.left) or
                      not is_approximately_equal(shape1.top, shape2.top) or
                      not is_approximately_equal(shape1.width, shape2.width) or
                      not is_approximately_equal(shape1.height, shape2.height)):
                    return 0

            if examine_modify_height:
                if not hasattr(shape1, "text") and not hasattr(shape2, "text") or shape1.shape_type == 5 and shape2.shape_type == 5:
                    if not is_approximately_equal(shape1.height, shape2.height):
                        return 0
                elif (not is_approximately_equal(shape1.left, shape2.left) or
                      not is_approximately_equal(shape1.top, shape2.top) or
                      not is_approximately_equal(shape1.width, shape2.width) or
                      not is_approximately_equal(shape1.height, shape2.height)):
                    return 0

            if shape1.shape_type == MSO_SHAPE_TYPE.TABLE:
                table1 = shape1.table
                table2 = shape2.table
                if len(table1.rows) != len(table2.rows) or len(table1.columns) != len(table2.columns):
                    return 0

                for row_idx in range(len(table1.rows)):
                    for col_idx in range(len(table1.columns)):
                        cell1 = table1.cell(row_idx, col_idx)
                        cell2 = table2.cell(row_idx, col_idx)

                        if len(cell1.text_frame.paragraphs) != len(cell2.text_frame.paragraphs):
                            return 0

                        for para1, para2 in zip(cell1.text_frame.paragraphs, cell2.text_frame.paragraphs):
                            runs1 = para1.runs
                            runs2 = para2.runs
                            if (para1.text or "").strip() == "" and (para2.text or "").strip() == "":
                                runs1 = nonempty_runs(para1)
                                runs2 = nonempty_runs(para2)

                            if len(runs1) != len(runs2) and examine_run_count:
                                return 0

                            for run1, run2 in zip(runs1, runs2):
                                if hasattr(run1.font.color, "rgb") and hasattr(run2.font.color, "rgb"):
                                    if run1.font.color.rgb != run2.font.color.rgb and examine_color_rgb:
                                        return 0
                                if run1.font.bold != run2.font.bold:
                                    if not ((run1.font.bold is None or run1.font.bold is False) and
                                           (run2.font.bold is None or run2.font.bold is False)):
                                        return 0
                                if run1.font.italic != run2.font.italic:
                                    if not ((run1.font.italic is None or run1.font.italic is False) and
                                           (run2.font.italic is None or run2.font.italic is False)):
                                        return 0
                                if run1.font.underline != run2.font.underline:
                                    if run1.font.underline is not None and run2.font.underline is not None:
                                        return 0
                                    if (run1.font.underline is None and run2.font.underline is True) or (run1.font.underline is True and run2.font.underline is None):
                                        return 0

            if hasattr(shape1, "text") and hasattr(shape2, "text"):
                if shape1.text.strip() != shape2.text.strip() and examine_text:
                    return 0

                if len(shape1.text_frame.paragraphs) != len(shape2.text_frame.paragraphs):
                    return 0

                for para1, para2 in zip(shape1.text_frame.paragraphs, shape2.text_frame.paragraphs):
                    if examine_alignment:
                        from pptx.enum.text import PP_ALIGN
                        align1 = para1.alignment
                        align2 = para2.alignment
                        if align1 is None:
                            align1 = PP_ALIGN.LEFT
                        if align2 is None:
                            align2 = PP_ALIGN.LEFT
                        if align1 != align2:
                            return 0

                    if para1.text != para2.text and examine_text:
                        return 0

                    if para1.level != para2.level and examine_indent:
                        n1 = 0 if para1.level is None else para1.level
                        n2 = 0 if para2.level is None else para2.level
                        if not (n1 in (0, 1) and n2 in (0, 1)):
                            return 0

                    runs1 = para1.runs
                    runs2 = para2.runs
                    if (para1.text or "").strip() == "" and (para2.text or "").strip() == "":
                        runs1 = nonempty_runs(para1)
                        runs2 = nonempty_runs(para2)

                    if len(runs1) != len(runs2) and examine_run_count:
                        return 0

                    for run1, run2 in zip(runs1, runs2):
                        if run1.font.name != run2.font.name and examine_font_name:
                            return 0
                        if run1.font.size != run2.font.size and examine_font_size:
                            return 0
                        if run1.font.bold != run2.font.bold and examine_font_bold:
                            if not ((run1.font.bold is None or run1.font.bold is False) and
                                   (run2.font.bold is None or run2.font.bold is False)):
                                return 0
                        if run1.font.italic != run2.font.italic and examine_font_italic:
                            if not ((run1.font.italic is None or run1.font.italic is False) and
                                   (run2.font.italic is None or run2.font.italic is False)):
                                return 0
                        if hasattr(run1.font.color, "rgb") and hasattr(run2.font.color, "rgb"):
                            if run1.font.color.rgb != run2.font.color.rgb and examine_color_rgb:
                                return 0
                        if run1.font.underline != run2.font.underline and examine_font_underline:
                            if run1.font.underline is not None and run2.font.underline is not None:
                                return 0
                            if (run1.font.underline is None and run2.font.underline is True) or (run1.font.underline is True and run2.font.underline is None):
                                return 0
                        if run1.font._element.attrib.get('strike', 'noStrike') != run2.font._element.attrib.get(
                                'strike', 'noStrike') and examine_strike_through:
                            return 0

                        def _extract_bullets(xml_data):
                            root = ET.fromstring(xml_data)
                            namespaces = {
                                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                            }
                            bullets = []
                            for paragraph in root.findall('.//a:p', namespaces):
                                pPr = paragraph.find('a:pPr', namespaces)
                                if pPr is not None:
                                    lvl = pPr.get('lvl')
                                    buChar = pPr.find('a:buChar', namespaces)
                                    char = buChar.get('char') if buChar is not None else "No Bullet"
                                    buClr = pPr.find('a:buClr/a:srgbClr', namespaces)
                                    color = buClr.get('val') if buClr is not None else "No Color"
                                else:
                                    lvl = "No Level"
                                    char = "No Bullet"
                                    color = "No Color"

                                text = "".join(t.text for t in paragraph.findall('.//a:t', namespaces))
                                if text.strip():
                                    bullets.append((lvl, char, text, color))
                            return bullets

                        def _compare_bullets_with_tolerance(bullets1, bullets2):
                            if len(bullets1) != len(bullets2):
                                return False
                            for (lvl1, char1, text1, color1), (lvl2, char2, text2, color2) in zip(bullets1, bullets2):
                                if text1 != text2:
                                    return False
                                if char1 != char2:
                                    return False
                                if char1 != "No Bullet" or char2 != "No Bullet":
                                    normalized_lvl1 = '0' if lvl1 is None else lvl1
                                    normalized_lvl2 = '0' if lvl2 is None else lvl2
                                    if normalized_lvl1 != normalized_lvl2:
                                        return False
                            return True

                        if examine_bullets:
                            try:
                                bullets1 = _extract_bullets(run1.part.blob.decode('utf-8'))
                                bullets2 = _extract_bullets(run2.part.blob.decode('utf-8'))
                                if not _compare_bullets_with_tolerance(bullets1, bullets2):
                                    return 0
                            except Exception:
                                pass

        if examine_alignment and len(text_shapes1) == len(text_shapes2):
            for tshape1, tshape2 in zip(text_shapes1, text_shapes2):
                if tshape1.text.strip() != tshape2.text.strip() and examine_text:
                    return 0
                if len(tshape1.text_frame.paragraphs) != len(tshape2.text_frame.paragraphs):
                    return 0
                for para1, para2 in zip(tshape1.text_frame.paragraphs, tshape2.text_frame.paragraphs):
                    from pptx.enum.text import PP_ALIGN
                    align1 = para1.alignment
                    align2 = para2.alignment
                    if align1 is None:
                        align1 = PP_ALIGN.LEFT
                    if align2 is None:
                        align2 = PP_ALIGN.LEFT
                    if align1 != align2:
                        return 0
        elif len(text_shapes1) != len(text_shapes2):
            return 0

    return 1


def check_slide_orientation_Portrait(pptx_path):
    presentation = Presentation(pptx_path)
    slide_height = presentation.slide_height
    slide_width = presentation.slide_width
    if slide_width < slide_height:
        return 1
    return 0


def evaluate_presentation_fill_to_rgb_distance(pptx_file, rules):
    rgb = rules["rgb"]
    try:
        original_rgb = rules["original_rgb"]
    except Exception:
        original_rgb = None

    def get_rgb_from_color(color):
        try:
            if hasattr(color, "rgb"):
                return color.rgb
            else:
                return None
        except Exception:
            return None

    def slide_fill_distance_to_rgb(_slide, _rgb, _original_rgb):
        fill = _slide.background.fill
        if fill.type == 1:
            color_rgb = get_rgb_from_color(fill.fore_color)
            if color_rgb is None:
                return 1
            r1, g1, b1 = color_rgb
            r2, g2, b2 = _rgb

            if _original_rgb is not None:
                r3, g3, b3 = _original_rgb
                if r1 == r3 and g1 == g3 and b1 == b3:
                    return 1

            return sqrt((r1 - r2) ** 2 + (g1 - g2) ** 2 + (b1 - b2) ** 2) / sqrt(255 ** 2 + 255 ** 2 + 255 ** 2)
        elif fill.type == 5:
            master_fill = _slide.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                color_rgb = get_rgb_from_color(master_fill.fore_color)
                if color_rgb is None:
                    return 1
                r1, g1, b1 = color_rgb
            else:
                return 1
            r2, g2, b2 = _rgb

            if _original_rgb is not None:
                r3, g3, b3 = _original_rgb
                if r1 == r3 and g1 == g3 and b1 == b3:
                    return 1

            return sqrt((r1 - r2) ** 2 + (g1 - g2) ** 2 + (b1 - b2) ** 2) / sqrt(255 ** 2 + 255 ** 2 + 255 ** 2)

        return 1

    prs = Presentation(pptx_file)
    similarity = 1 - sum(slide_fill_distance_to_rgb(slide, rgb, original_rgb) for slide in prs.slides) / len(prs.slides)
    return similarity


def check_transition(pptx_file, rules):
    slide_idx = rules['slide_idx']
    transition_type = rules['transition_type']

    with zipfile.ZipFile(pptx_file, 'r') as zip_ref:
        slide_name = 'ppt/slides/slide{}.xml'.format(slide_idx + 1)
        try:
            zip_ref.getinfo(slide_name)
        except KeyError:
            return 0.

        with zip_ref.open(slide_name) as slide_file:
            tree = ET.parse(slide_file)
            root = tree.getroot()

            namespaces = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            }

            transition = root.find('.//p:transition', namespaces)
            if transition is not None:
                dissolve = transition.find('.//p:{}'.format(transition_type), namespaces)
                if dissolve is not None:
                    return 1.
                else:
                    return 0.
            else:
                return 0.


def check_page_number_colors(pptx_file, rules):
    color = rules["color"]

    def parse_rgb(rgb_str):
        if rgb_str is None:
            return None
        rgb_str = rgb_str.lstrip('#')
        if len(rgb_str) != 6:
            return None
        try:
            r = int(rgb_str[0:2], 16)
            g = int(rgb_str[2:4], 16)
            b = int(rgb_str[4:6], 16)
            return (r, g, b)
        except ValueError:
            return None

    def is_red(rgb_tuple, threshold=50):
        if rgb_tuple is None:
            return False
        r, g, b = rgb_tuple
        return r > g + threshold and r > b + threshold

    def is_blue(rgb_tuple, threshold=50):
        if rgb_tuple is None:
            return False
        r, g, b = rgb_tuple
        return b > g + threshold and b > r + threshold

    def is_green(rgb_tuple, threshold=50):
        if rgb_tuple is None:
            return False
        r, g, b = rgb_tuple
        return g > r + threshold and g > b + threshold

    def is_black(rgb_tuple, threshold=50):
        if rgb_tuple is None:
            return False
        r, g, b = rgb_tuple
        return r < threshold and g < threshold and b < threshold

    with zipfile.ZipFile(pptx_file, 'r') as zip_ref:
        slide_master_name = 'ppt/slideMasters/slideMaster1.xml'
        with zip_ref.open(slide_master_name) as slide_master_file:
            tree = ET.parse(slide_master_file)
            root = tree.getroot()

            namespaces = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            }

            slide_number_ph = root.find('.//p:ph[@type="sldNum"]', namespaces)
            slides_color_val = None

            if slide_number_ph is not None:
                color_elem = slide_number_ph.find('.//a:solidFill//a:srgbClr', namespaces)
                if color_elem is not None:
                    slides_color_val = color_elem.get('val')

            if slides_color_val is None:
                all_ph = root.findall('.//p:ph', namespaces)
                for ph in all_ph:
                    ph_type = ph.get('type')
                    if ph_type == 'sldNum' or ph_type == 'ftr' or ph_type == 'dt':
                        color_elem = ph.find('.//a:solidFill//a:srgbClr', namespaces)
                        if color_elem is not None:
                            slides_color_val = color_elem.get('val')
                            break

            if slides_color_val is None:
                text_runs = root.findall('.//a:rPr//a:solidFill//a:srgbClr', namespaces)
                if text_runs:
                    for color_elem in reversed(text_runs):
                        color_val = color_elem.get('val')
                        if color_val and color_val != '000000':
                            slides_color_val = color_val
                            break

            if slides_color_val is None:
                color_elems = root.findall('.//a:solidFill//a:srgbClr', namespaces)
                for color_elem in reversed(color_elems):
                    color_val = color_elem.get('val')
                    if color_val and color_val != '000000':
                        slides_color_val = color_val
                        break
                if slides_color_val is None and color_elems:
                    slides_color_val = color_elems[-1].get('val')

    if slides_color_val is None:
        return 0

    rgb_tuple = parse_rgb(slides_color_val)
    if rgb_tuple is None:
        return 0

    if color == "red" and not is_red(rgb_tuple):
        return 0
    elif color == "blue" and not is_blue(rgb_tuple):
        return 0
    elif color == "green" and not is_green(rgb_tuple):
        return 0
    elif color == "black" and not is_black(rgb_tuple):
        return 0

    return 1


def check_auto_saving_time(pptx_file, rules):
    minutes = rules["minutes"]

    try:
        tree = ET.parse(pptx_file)
        root = tree.getroot()
        autosave_time = None
        for item in root.findall(".//item"):
            path = item.get('{http://openoffice.org/2001/registry}path')
            if path == "/org.openoffice.Office.Common/Save/Document":
                for prop in item.findall(".//prop"):
                    name = prop.get('{http://openoffice.org/2001/registry}name')
                    if name == "AutoSaveTimeIntervall":
                        autosave_time = prop.find(".//value").text
                        break
        if autosave_time is None:
            return 0
        else:
            autosave_time = int(autosave_time)
            if autosave_time == minutes:
                return 1
            else:
                return 0
    except Exception as e:
        logger.error(f"Error: {e}")
        return 0


def compare_images(image1_path, image2_path, **options):
    if not image1_path or not image2_path:
        return 0

    _require_optional_dependencies(Image=Image, np=np, ssim=ssim)

    base_score = options.get("reference_base_result", None)
    image1 = Image.open(image1_path).convert('L')
    image2 = Image.open(image2_path).convert('L')
    image1_size = image1.size
    image2_size = image2.size
    new_size = min(image1_size, image2_size)
    image1 = image1.resize(new_size, Image.Resampling.LANCZOS)
    image2 = image2.resize(new_size, Image.Resampling.LANCZOS)
    image1_array = np.array(image1)
    image2_array = np.array(image2)
    similarity_index = ssim(image1_array, image2_array)

    epsilon = 0.01
    if base_score is not None:
        if similarity_index >= base_score + epsilon:
            return (similarity_index - base_score) / (1 - base_score)
        else:
            return 0
    else:
        return similarity_index


def compare_audios(audio_path_1, audio_path_2):
    if not audio_path_1 or not audio_path_2:
        return 0

    _require_optional_dependencies(librosa=librosa, np=np, fastdtw=fastdtw, cosine=cosine)

    y1, y2 = None, None
    try:
        y1, sr1 = librosa.load(audio_path_1)
    except Exception:
        logger.warning(f"Could not load audio from {os.path.basename(audio_path_1)}. It might be empty or corrupt.")

    try:
        y2, sr2 = librosa.load(audio_path_2)
    except Exception:
        logger.warning(f"Could not load audio from {os.path.basename(audio_path_2)}. It might be empty or corrupt.")

    is_y1_bad = (y1 is None) or (y1.shape[0] == 0)
    is_y2_bad = (y2 is None) or (y2.shape[0] == 0)

    if is_y1_bad and is_y2_bad:
        logger.info("Both audio files are empty or corrupt. Considering them perfectly similar.")
        return 1.0

    if is_y1_bad or is_y2_bad:
        logger.warning("One audio file is empty/corrupt, the other is not. Similarity is 0.")
        return 0.0

    try:
        mfcc1 = librosa.feature.mfcc(y=y1, sr=sr1)
        mfcc2 = librosa.feature.mfcc(y=y2, sr=sr2)
    except Exception as e:
        logger.error(f"Error during MFCC extraction: {e}")
        return 0.0

    mfcc1 = librosa.util.normalize(mfcc1, axis=1)
    mfcc2 = librosa.util.normalize(mfcc2, axis=1)

    dist_func = lambda x, y: cosine(x, y)
    distance, path = fastdtw(mfcc1.T, mfcc2.T, dist=dist_func)

    if len(path) == 0:
        normalized_distance = np.inf
    else:
        normalized_distance = distance / len(path)

    similarity = np.exp(-normalized_distance)
    return similarity


def get_audio_in_slide_local(config: Dict[str, str]):
    ppt_file_path = localize_path(config["ppt_file_path"])
    slide_index = int(config["slide_index"])
    audio_file_path = None

    with zipfile.ZipFile(ppt_file_path, 'r') as myzip:
        slide_rels_file = 'ppt/slides/_rels/slide{}.xml.rels'.format(slide_index + 1)
        if slide_rels_file in myzip.namelist():
            with myzip.open(slide_rels_file) as f:
                tree = ET.parse(f)
                root = tree.getroot()
                namespaces = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
                for rel in root.findall('r:Relationship', namespaces):
                    if 'audio' in rel.attrib['Type']:
                        target = rel.attrib['Target']
                        if target.startswith('..'):
                            audio_file_path = os.path.normpath(os.path.join('ppt/slides', target))
                            audio_file_path = audio_file_path.replace('\\', '/')
                            tmpdirname = os.path.dirname(str(ppt_file_path))
                            myzip.extract(audio_file_path, tmpdirname)
                            audio_file_path = os.path.join(tmpdirname, audio_file_path)
                            return audio_file_path
                        else:
                            assert target.startswith("file://"), target
                            audio_file_path = target[7:]
                        break
    if audio_file_path is None:
        return None
    return str(localize_path(audio_file_path))


def resolve_expected(spec):
    if spec is None:
        return None
    if spec["type"] == "cloud_file":
        return str(fetch_cloud_file(spec["url"], spec["dest"]))
    if spec["type"] == "rule":
        return spec["rules"]
    raise ValueError(f"Unsupported expected spec: {spec}")


def resolve_result_state(state):
    if isinstance(state, str):
        return str(localize_path(state))
    if state["type"] == "audio_in_slide":
        return get_audio_in_slide_local(state)
    raise ValueError(f"Unsupported result state: {state}")


def evaluate_single(func_name, result_state, expected_state=None, options=None):
    options = options or {}
    metric = {
        "check_presenter_console_disable": check_presenter_console_disable,
        "check_image_stretch_and_center": check_image_stretch_and_center,
        "compare_pptx_files": compare_pptx_files,
        "check_slide_orientation_Portrait": check_slide_orientation_Portrait,
        "evaluate_presentation_fill_to_rgb_distance": evaluate_presentation_fill_to_rgb_distance,
        "check_transition": check_transition,
        "check_page_number_colors": check_page_number_colors,
        "check_auto_saving_time": check_auto_saving_time,
        "compare_images": compare_images,
        "compare_audios": compare_audios,
    }[func_name]

    try:
        if expected_state is None:
            return float(metric(result_state, **options))
        return float(metric(result_state, expected_state, **options))
    except Exception as exc:
        logger.error("Evaluation failed for %s: %s", func_name, exc)
        return 0.0


def run_evaluation():
    metric_spec = TASK_SPEC["metric_spec"]
    func = metric_spec["func"]

    if isinstance(func, list):
        conj = metric_spec.get("conj", "and")
        results = []
        for index, func_name in enumerate(func):
            result_state = resolve_result_state(metric_spec["result"][index])
            expected_state = resolve_expected(metric_spec["expected"][index]) if "expected" in metric_spec else None
            options = metric_spec.get("options", [{}] * len(func))[index]
            metric_value = evaluate_single(func_name, result_state, expected_state, options)
            if conj == "and" and metric_value == 0.0:
                return 0.0
            if conj == "or" and metric_value == 1.0:
                return 1.0
            results.append(metric_value)
        return sum(results) / len(results) if conj == "and" else max(results)

    result_state = resolve_result_state(metric_spec["result"])
    expected_state = resolve_expected(metric_spec["expected"]) if "expected" in metric_spec else None
    return evaluate_single(func, result_state, expected_state, metric_spec.get("options", {}))


def run_sanity():
    with tempfile.TemporaryDirectory(prefix=f"{TASK_SPEC['slug']}-") as temp_dir:
        app_dir = Path(temp_dir) / "app"
        agent_home = Path(temp_dir) / "home" / "agent"
        previous_app_dir = os.environ.get("APP_DIR")
        previous_agent_home = os.environ.get("AGENT_HOME")
        os.environ["APP_DIR"] = str(app_dir)
        os.environ["AGENT_HOME"] = str(agent_home)
        try:
            materialize_inputs(app_dir)
            fail_score = run_evaluation()
            run_oracle(app_dir, agent_home)
            pass_score = run_evaluation()
        finally:
            if previous_app_dir is None:
                os.environ.pop("APP_DIR", None)
            else:
                os.environ["APP_DIR"] = previous_app_dir
            if previous_agent_home is None:
                os.environ.pop("AGENT_HOME", None)
            else:
                os.environ["AGENT_HOME"] = previous_agent_home
    return {"fail_score": fail_score, "pass_score": pass_score}


if __name__ == "__main__":
    import sys

    if len(sys.argv) > 1 and sys.argv[1] == "--sanity":
        print(json.dumps(run_sanity(), sort_keys=True))
    else:
        print(run_evaluation())
